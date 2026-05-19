from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

def add_slide(title_text, content_list):
    """Helper function to create a slide with a title and bullet points."""
    slide_layout = prs.slide_layouts[1] # 1 is Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, item in enumerate(content_list):
        if i == 0:
            tf.text = item
        else:
            p = tf.add_paragraph()
            p.text = item

# SLIDE 1: Title
slide_layout = prs.slide_layouts[0] # 0 is Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ULTRA PRO ACCESS CORE AI"
subtitle.text = "Biometric Authentication & Reactive Hardware Actuation\n\nDeveloped by: Yassine Mokni & Hadil Dhaya"

# SLIDE 2: Problem
add_slide("The Access Control Crisis", [
    "Legacy Vulnerabilities: Traditional systems relying on keycards and PINs are fundamentally flawed. Hardware tokens can be cloned, shared, or lost.",
    "The Latency Barrier: Modern cloud-based AI solutions introduce critical network latency. Access control requires localized, sub-millisecond execution."
])

# SLIDE 3: Solution
add_slide("The Unified Solution", [
    "Safe Exit Pro bridges the gap between Deep Learning and physical hardware.",
    "Localized Inference: AI runs directly on Apple Silicon without cloud delays.",
    "Zero-Delay Actuators: Deterministic serial communication achieves instant hardware actuation based on biometric verification."
])

# SLIDE 4: Data Scale
add_slide("Unprecedented Data Scale: 1.1 Million Images", [
    "To ensure 99%+ validation accuracy, the model was trained on a massive global dataset.",
    "Dataset size: 1,100,000 images.",
    "Coverage: Diverse lighting conditions, angles, and facial features for enterprise-grade generalization."
])

# SLIDE 5: Strategy
add_slide("Phase 1: Subset Benchmarking", [
    "Rapid Prototyping: Before full-scale deployment, we isolated a strategic subset of 10,000 images.",
    "Optimization: This allowed us to benchmark hyper-parameters and optimize the Metal Performance Shaders (MPS) pipeline.",
    "Hardware Sync: Validated the serial bus response time in a controlled environment."
])

# SLIDE 6: Results
add_slide("AI Performance: Accuracy Curve", [
    "Subset Phase (10,000 images): Achieved 88.2% Accuracy.",
    "Full Deployment (1.1 Million images): Scaled up to achieve 99.4% Accuracy.",
    "Result: Pushed biometric accuracy to enterprise standards while maintaining inference speed."
])

# SLIDE 7: Hardware
add_slide("The Hardware Node: Arduino Actuation", [
    "Deterministic Firmware: C++ Arduino logic designed strictly without delay() to prevent buffer backlog.",
    "Serial Bus Bridge: Single-byte command protocol ('G', 'D', 'O') for minimal overhead.",
    "Fail-Safe Protocol: Automatic port discovery and clean exit flushing to ensure hardware safely powers down."
])

# SLIDE 8: Stabilization
add_slide("Security Logic: Stabilization", [
    "Temporal Confirmation: Implements a 5-Frame Temporal Lock to eliminate flickering. The actuator triggers only after 5 consecutive matches.",
    "Dominant Face Selection: In multi-person scenes, the algorithm isolates the closest face (by spatial area) for authoritative decision making."
])

# SLIDE 9: Features
add_slide("Production Features", [
    "MPS Acceleration: Direct mapping to Apple Silicon GPUs for extreme low-latency inference.",
    "Multi-Modal Feedback: Simultaneous LED triggering and macOS voice feedback.",
    "Rapid Enrollment: One-click identity capture (150 frames) with automatic mean embedding calculation.",
    "Safe Exit Protocol: Hardware-level serial flushing prevents ghost signals."
])

# SLIDE 10: Future
add_slide("System Expansion & Roadmap", [
    "Magnetic Locks & Relays: Upgrading from indicator LEDs to high-voltage magnetic door locks for full facility integration.",
    "Anti-Spoofing AI: Deployment of liveness detection (blink/depth sensing) to prevent photo spoofing.",
    "Encrypted Logs: Cryptographically secure access exports for enterprise compliance."
])

# SLIDE 11: Q&A
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Questions?"
subtitle.text = "Ultra Pro Access Core AI\nTeam: Yassine Mokni & Hadil Dhaya\nGitHub.com/Yaxin02"

# Save the presentation
filename = "Ultra_Pro_Access_Core_AI_Presentation.pptx"
prs.save(filename)
print(f"✅ Success! PowerPoint file '{filename}' has been created on your Mac.")